#!/usr/bin/env python3
"""Post-process pandoc pptx to fix multi-image slides.

Pandoc 2.x splits multiple images under one heading into separate slides,
and only embeds the first image — subsequent ones become text-only orphans.

This script:
  1. Parses slides.md to learn which slide titles map to which images
  2. Removes orphan slides (pandoc-generated splits with no real content)
  3. Lays out all intended images on the correct parent slide

Usage:
    python scripts/postprocess_slides.py slides.pptx slides.md [--resource-path DIR]

If --resource-path is not given, images are resolved relative to slides.md's directory.
"""
from __future__ import annotations

import io
import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Emu
except ImportError:
    print("python-pptx is required: pip install python-pptx", file=sys.stderr)
    sys.exit(1)


# ---------------------------------------------------------------------------
# Parse slides.md to extract slide → images mapping
# ---------------------------------------------------------------------------

def parse_yaml_notes(md_path: Path) -> str:
    """Extract the top-level 'notes:' value from YAML frontmatter, if present."""
    text = md_path.read_text(encoding="utf-8")
    if not text.startswith("---\n"):
        return ""
    end = text.find("\n---\n", 4)
    if end == -1:
        return ""
    fm_text = text[4:end]
    # Simple key extraction — avoid pulling in a full YAML parser dependency
    match = re.search(r"(?m)^notes:\s*\|\n((?:[ \t]+.+\n?)+)", fm_text)
    if match:
        lines = match.group(1).splitlines()
        return "\n".join(line.strip() for line in lines).strip()
    match = re.search(r"(?m)^notes:\s*['\"]?(.*?)['\"]?\s*$", fm_text)
    if match:
        return match.group(1).strip()
    return ""


def parse_slide_images(md_path: Path, resource_path: Path) -> list[dict]:
    """Parse slides.md and return list of {title, images, has_text} per slide."""
    text = md_path.read_text(encoding="utf-8")

    # Strip YAML frontmatter
    if text.startswith("---\n"):
        end = text.find("\n---\n", 4)
        if end != -1:
            text = text[end + 5:]

    # Collect all image alt texts for split detection
    all_alts: set[str] = set()

    # Split on ## headings
    slides = []
    matches = list(re.finditer(r"(?m)^##\s+(.+)$", text))
    for i, m in enumerate(matches):
        title = m.group(1).strip()
        start = m.end()
        end = matches[i + 1].start() if i + 1 < len(matches) else len(text)
        block = text[start:end]

        # Extract image references
        images = []
        for img_match in re.finditer(r"!\[([^\]]*)\]\(([^)]+)\)", block):
            alt = img_match.group(1)
            src = img_match.group(2)
            all_alts.add(normalize(alt))
            full_path = resource_path / src
            if full_path.exists():
                images.append({"alt": alt, "path": full_path})

        # Check if block has bullet text (not just images and notes)
        # Strip images, notes blocks, and blank lines — see if bullets remain
        stripped = re.sub(r"!\[[^\]]*\]\([^)]+\)", "", block)
        stripped = re.sub(r":::.*?:::", "", stripped, flags=re.DOTALL)
        has_text = bool(re.search(r"(?m)^(- .+|\|.+\|)", stripped))

        slides.append({"title": title, "images": images, "has_text": has_text})

    return slides, all_alts


# ---------------------------------------------------------------------------
# PPTX helpers
# ---------------------------------------------------------------------------

def get_slide_title(slide) -> str:
    """Extract title text from a slide."""
    for shape in slide.shapes:
        try:
            if shape.placeholder_format is not None and shape.placeholder_format.idx == 0:
                return shape.text_frame.text.strip()
        except (ValueError, AttributeError):
            pass
    # Fallback: first text shape
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            return shape.text_frame.text.strip()
    return ""


def get_slide_images(slide) -> list:
    """Return image shapes on a slide."""
    return [s for s in slide.shapes if hasattr(s, "image")]


def get_table_frames(slide) -> list:
    """Return ALL shapes (including body placeholders) that contain a table.

    Pandoc renders markdown tables inside the body placeholder (ph_idx=1) as a
    graphicFrame.  We must NOT skip placeholder shapes here, or we will miss
    the table that pandoc placed on the parent slide.
    """
    from pptx.oxml.ns import qn
    frames = []
    for shape in slide.shapes:
        # Skip title placeholders (idx==0); include body (idx==1) and standalone
        try:
            phf = shape.placeholder_format
            if phf is not None and phf.idx == 0:
                continue  # title placeholder — never a table
        except (ValueError, AttributeError):
            pass
        if shape._element.find(f".//{qn('a:tbl')}") is not None:
            frames.append(shape)
    return frames


def get_text_shapes(slide) -> list:
    """Return non-placeholder text shapes that have visible content.

    After _clone_text_as_standalone() the cloned bullet body becomes a
    non-placeholder sp element.  This function finds those shapes so
    layout_tables_on_slide() can position them.
    """
    shapes = []
    for shape in slide.shapes:
        try:
            if shape.placeholder_format is not None:
                continue  # skip title / body placeholders
        except (ValueError, AttributeError):
            pass
        if shape.has_text_frame and shape.text_frame.text.strip():
            shapes.append(shape)
    return shapes


def _next_shape_id(spTree) -> int:
    """Return a shape ID that is not yet used in spTree."""
    from pptx.oxml.ns import qn
    max_id = 0
    for sp in spTree:
        for cNvPr in sp.findall(f".//{qn('p:cNvPr')}"):
            try:
                max_id = max(max_id, int(cNvPr.get("id", "0")))
            except ValueError:
                pass
    return max_id + 1


def _transfer_table_as_standalone(parent_slide, orphan_table_shape) -> None:
    """Copy a table from an orphan slide onto the parent as a NEW standalone graphicFrame.

    The orphan's table lives inside a body placeholder (ph_idx=1).  A plain
    deepcopy would add a second ph_idx=1 element — invalid in pptx.  Instead
    we create a fresh non-placeholder graphicFrame and embed only the a:tbl.
    """
    from copy import deepcopy
    from lxml import etree
    from pptx.oxml.ns import qn

    NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
    NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    TABLE_URI = "http://schemas.openxmlformats.org/drawingml/2006/table"

    tbl_elem = orphan_table_shape._element.find(f".//{qn('a:tbl')}")
    if tbl_elem is None:
        return

    # Read position/size from the orphan shape's xfrm
    xfrm_elem = orphan_table_shape._element.find(qn("p:xfrm"))
    x, y, cx, cy = "457200", "274638", "8229600", "4525963"
    if xfrm_elem is not None:
        off = xfrm_elem.find(qn("a:off"))
        ext = xfrm_elem.find(qn("a:ext"))
        if off is not None:
            x, y = off.get("x", x), off.get("y", y)
        if ext is not None:
            cx, cy = ext.get("cx", cx), ext.get("cy", cy)

    spTree = parent_slide._element.find(qn("p:cSld")).find(qn("p:spTree"))
    new_id = _next_shape_id(spTree)

    gf_xml = (
        f'<p:graphicFrame xmlns:p="{NS_P}" xmlns:a="{NS_A}">'
        f'<p:nvGraphicFramePr>'
        f'<p:cNvPr id="{new_id}" name="Table {new_id}"/>'
        f'<p:cNvGraphicFramePr><a:graphicFrameLocks noGrp="1"/></p:cNvGraphicFramePr>'
        f'<p:nvPr/>'
        f'</p:nvGraphicFramePr>'
        f'<p:xfrm><a:off x="{x}" y="{y}"/><a:ext cx="{cx}" cy="{cy}"/></p:xfrm>'
        f'<a:graphic><a:graphicData uri="{TABLE_URI}"/></a:graphic>'
        f'</p:graphicFrame>'
    )
    gf_elem = etree.fromstring(gf_xml)
    graphicData = gf_elem.find(f"{{{NS_A}}}graphic/{{{NS_A}}}graphicData")
    graphicData.append(deepcopy(tbl_elem))
    spTree.append(gf_elem)


def _clone_text_as_standalone(parent_slide, orphan_body) -> None:
    """Clone an orphan's body placeholder as a non-placeholder text shape.

    When the parent's body placeholder holds a TABLE (has no text_frame), we
    cannot merge bullets into it.  Instead we deepcopy the orphan's body sp
    element, strip the <p:ph> marker to make it standalone, and append it to
    the parent's spTree.  layout_tables_on_slide() will then position it in
    the text zone below the table(s).
    """
    from copy import deepcopy
    from pptx.oxml.ns import qn

    if not orphan_body.has_text_frame:
        return
    if not orphan_body.text_frame.text.strip():
        return

    spTree = parent_slide._element.find(qn("p:cSld")).find(qn("p:spTree"))
    new_id = _next_shape_id(spTree)

    sp_clone = deepcopy(orphan_body._element)

    # Remove <p:ph> so PowerPoint treats this as a standalone shape
    for nvPr in sp_clone.findall(f".//{qn('p:nvPr')}"):
        for ph in list(nvPr.findall(qn("p:ph"))):
            nvPr.remove(ph)

    # Assign a new shape ID
    for cNvPr in sp_clone.findall(f".//{qn('p:cNvPr')}"):
        cNvPr.set("id", str(new_id))
        break

    # Explicitly set font size (1600 = 16 pt) so theme inheritance isn't needed
    for rPr in sp_clone.findall(f".//{qn('a:rPr')}"):
        if not rPr.get("sz"):
            rPr.set("sz", "1600")

    spTree.append(sp_clone)


def has_real_content(slide) -> bool:
    """Check if slide has meaningful content beyond just a caption."""
    images = get_slide_images(slide)
    if images:
        return True
    for s in slide.shapes:
        try:
            if s.placeholder_format is None:
                continue
            idx = s.placeholder_format.idx
            if idx == 0 and s.has_text_frame:
                # Title with more than 4 words counts
                if len(s.text_frame.text.strip().split()) > 4:
                    return True
            elif idx == 1:
                # Any body content counts — text or table/graphic
                if s.has_text_frame and s.text_frame.text.strip():
                    return True
                if not s.has_text_frame:
                    # Table or graphic placeholder — has real content
                    return True
        except (ValueError, AttributeError):
            pass
    return False


def normalize(s: str) -> str:
    """Normalize title for fuzzy matching. Replaces smart quotes and collapses whitespace."""
    s = s.strip().lower()
    # Normalize smart quotes to ASCII
    s = s.replace("\u2018", "'").replace("\u2019", "'")  # ' '
    s = s.replace("\u201c", '"').replace("\u201d", '"')  # " "
    s = s.replace("\u2013", "-").replace("\u2014", "-")  # – —
    s = re.sub(r"\s+", " ", s)
    return s[:80]


def get_body_placeholder(slide):
    """Return the body placeholder (idx=1) if it exists."""
    for shape in slide.shapes:
        try:
            if shape.placeholder_format is not None and shape.placeholder_format.idx == 1:
                return shape
        except (ValueError, AttributeError):
            pass
    return None


def slide_has_text_content(slide) -> bool:
    """Check if slide has real text bullets or table in the body placeholder."""
    body = get_body_placeholder(slide)
    if body is None:
        return False
    if not body.has_text_frame:
        # Table or graphic placeholder — counts as text content for layout purposes
        return True
    text = body.text_frame.text.strip()
    # Has at least 2 non-empty paragraphs (real bullets, not just image alt text)
    lines = [l for l in text.split("\n") if l.strip()]
    return len(lines) >= 2


def layout_images_with_text(slide, image_paths: list[Path], slide_w, slide_h):
    """Layout images in the TOP area and body text in the BOTTOM area.

    Used when a slide has both images and bullet text content.
    Images get ~58% of the content height (top), text gets ~37% (bottom).
    """
    from pptx.util import Inches, Emu

    n = len(image_paths)
    if n == 0:
        return

    top = Inches(1.3)
    bottom_margin = Inches(0.2)
    side_margin = Inches(0.4)
    gap_between = Inches(0.15)   # gap between image zone and text zone
    img_gap = Inches(0.2)        # gap between images when side-by-side

    avail_h = slide_h - top - bottom_margin
    avail_w = slide_w - 2 * side_margin

    # Split content area: top 58% images, bottom 37% text
    img_zone_h = int(avail_h * 0.58)
    text_zone_h = int(avail_h * 0.37)
    text_top = int(top + img_zone_h + gap_between)

    # Resize body placeholder to bottom zone
    body = get_body_placeholder(slide)
    if body is not None:
        body.left = int(side_margin)
        body.top = int(text_top)
        body.width = int(avail_w)
        body.height = int(text_zone_h)

    # Layout images in the top zone (2×2 grid for 4 images)
    if n == 1:
        cols, rows = 1, 1
    elif n <= 3:
        cols, rows = n, 1
    elif n <= 4:
        cols, rows = 2, 2
    else:
        cols = 3
        rows = (n + cols - 1) // cols

    cell_w = int((avail_w - img_gap * (cols - 1)) / cols)
    cell_h = int((img_zone_h - img_gap * (rows - 1)) / rows)

    for i, img_path in enumerate(image_paths):
        col = i % cols
        row = i // cols
        cell_x = int(side_margin + col * (cell_w + img_gap))
        cell_y = int(top + row * (cell_h + img_gap))

        pic = slide.shapes.add_picture(str(img_path), cell_x, cell_y)
        native_w = pic.width
        native_h = pic.height

        if native_w > 0 and native_h > 0:
            scale = min(cell_w / native_w, cell_h / native_h)
            new_w = int(native_w * scale)
            new_h = int(native_h * scale)
        else:
            new_w, new_h = cell_w, cell_h

        pic.width = new_w
        pic.height = new_h
        # Center horizontally and vertically within cell
        pic.left = cell_x + (cell_w - new_w) // 2
        pic.top = cell_y + (cell_h - new_h) // 2


def layout_images_on_slide(slide, image_paths: list[Path], slide_w, slide_h):
    """Add images to a slide in a responsive grid, preserving aspect ratio."""
    n = len(image_paths)
    if n == 0:
        return

    # Content area below header bar
    top = Inches(1.25)
    bottom_margin = Inches(0.25)
    left = Inches(0.4)
    right_margin = Inches(0.4)
    gap = Inches(0.2)

    avail_w = slide_w - left - right_margin
    avail_h = slide_h - top - bottom_margin

    if n == 1:
        cols, rows = 1, 1
    elif n == 2:
        cols, rows = 2, 1
    elif n == 3:
        cols, rows = 3, 1
    elif n <= 4:
        cols, rows = 2, 2
    else:
        cols = 3
        rows = (n + cols - 1) // cols

    cell_w = int((avail_w - gap * (cols - 1)) / cols)
    cell_h = int((avail_h - gap * (rows - 1)) / rows)

    for i, img_path in enumerate(image_paths):
        col = i % cols
        row = i // cols
        cell_x = int(left + col * (cell_w + gap))
        cell_y = int(top + row * (cell_h + gap))

        # Add image without specifying dimensions to get native size
        pic = slide.shapes.add_picture(str(img_path), cell_x, cell_y)

        # Read native size (the actual image dimensions in EMU)
        native_w = pic.width
        native_h = pic.height

        # Fit within cell while preserving aspect ratio
        if native_w > 0 and native_h > 0:
            scale = min(cell_w / native_w, cell_h / native_h)
            new_w = int(native_w * scale)
            new_h = int(native_h * scale)
        else:
            new_w, new_h = cell_w, cell_h

        pic.width = new_w
        pic.height = new_h

        # Center horizontally and vertically within the cell
        pic.left = cell_x + (cell_w - new_w) // 2
        pic.top = cell_y + (cell_h - new_h) // 2


def layout_stacked_tables(slide, slide_w: int, slide_h: int):
    """Reposition multiple table graphicFrames to stack vertically in the content area.

    Called after orphan table slides have been merged into their parent.
    Does nothing if the slide has fewer than 2 standalone table shapes.
    """
    from pptx.util import Inches
    from pptx.oxml.ns import qn

    table_frames = get_table_frames(slide)
    n = len(table_frames)
    if n < 2:
        return

    top = Inches(1.3)
    bottom_margin = Inches(0.2)
    side_margin = Inches(0.4)
    gap = Inches(0.12)

    avail_h = slide_h - top - bottom_margin
    avail_w = slide_w - 2 * side_margin

    cell_h = int((avail_h - gap * (n - 1)) / n)

    for i, shape in enumerate(table_frames):
        cell_y = int(top + i * (cell_h + gap))
        xfrm = shape._element.find(qn("p:xfrm"))
        if xfrm is None:
            continue
        off = xfrm.find(qn("a:off"))
        ext = xfrm.find(qn("a:ext"))
        if off is not None:
            off.set("x", str(int(side_margin)))
            off.set("y", str(cell_y))
        if ext is not None:
            ext.set("cx", str(int(avail_w)))
            ext.set("cy", str(cell_h))


# ---------------------------------------------------------------------------
# Table styling
# ---------------------------------------------------------------------------

def _hex_color(rgb: tuple) -> str:
    """Convert (R, G, B) tuple to 6-char hex string for OOXML."""
    return "%02X%02X%02X" % rgb


def _set_cell_fill(tc_elem, rgb: tuple):
    """Apply solid fill to a table cell XML element."""
    from lxml import etree
    from pptx.oxml.ns import qn as _qn

    tcPr = tc_elem.find(_qn("a:tcPr"))
    if tcPr is None:
        tcPr = etree.SubElement(tc_elem, _qn("a:tcPr"))

    for tag in [_qn("a:noFill"), _qn("a:solidFill"), _qn("a:gradFill"),
                _qn("a:blipFill"), _qn("a:pattFill"), _qn("a:grpFill")]:
        for child in list(tcPr.findall(tag)):
            tcPr.remove(child)

    solidFill = etree.SubElement(tcPr, _qn("a:solidFill"))
    srgbClr = etree.SubElement(solidFill, _qn("a:srgbClr"))
    srgbClr.set("val", _hex_color(rgb))


def _style_cell_runs(tc_elem, text_rgb: tuple, bold: bool):
    """Set text color and bold on every run in a table cell."""
    from lxml import etree
    from pptx.oxml.ns import qn as _qn

    for r_elem in tc_elem.findall(".//" + _qn("a:r")):
        rPr = r_elem.find(_qn("a:rPr"))
        if rPr is None:
            rPr = etree.Element(_qn("a:rPr"))
            r_elem.insert(0, rPr)

        rPr.set("b", "1" if bold else "0")

        for tag in [_qn("a:solidFill"), _qn("a:gradFill"), _qn("a:noFill"),
                    _qn("a:schemeClr"), _qn("a:sysClr")]:
            for child in list(rPr.findall(tag)):
                rPr.remove(child)

        solidFill = etree.SubElement(rPr, _qn("a:solidFill"))
        srgbClr = etree.SubElement(solidFill, _qn("a:srgbClr"))
        srgbClr.set("val", _hex_color(text_rgb))


def style_tables(prs):
    """Apply navy header / alternating-row styling to all tables in the deck.

    Header row → navy fill (#1B2A4A), white bold text
    Odd rows   → light blue fill (#E8F0FA), dark text
    Even rows  → white fill, dark text
    """
    from pptx.oxml.ns import qn as _qn

    HEADER_FILL = (0x1B, 0x2A, 0x4A)   # navy
    ODD_FILL    = (0xE8, 0xF0, 0xFA)   # light blue
    EVEN_FILL   = (0xFF, 0xFF, 0xFF)   # white
    HEADER_TEXT = (0xFF, 0xFF, 0xFF)   # white
    BODY_TEXT   = (0x1A, 0x1A, 0x1A)  # near-black

    for slide in prs.slides:
        for shape in slide.shapes:
            tbl = shape._element.find(".//" + _qn("a:tbl"))
            if tbl is None:
                continue

            for row_idx, tr in enumerate(tbl.findall(_qn("a:tr"))):
                if row_idx == 0:
                    fill, text, bold = HEADER_FILL, HEADER_TEXT, True
                elif row_idx % 2 == 1:
                    fill, text, bold = ODD_FILL, BODY_TEXT, False
                else:
                    fill, text, bold = EVEN_FILL, BODY_TEXT, False

                for tc in tr.findall(_qn("a:tc")):
                    _set_cell_fill(tc, fill)
                    _style_cell_runs(tc, text, bold)


# ---------------------------------------------------------------------------
# Main postprocess logic
# ---------------------------------------------------------------------------

def postprocess(pptx_path: Path, md_path: Path, resource_path: Path):
    slide_specs, all_image_alts = parse_slide_images(md_path, resource_path)
    title_notes = parse_yaml_notes(md_path)
    prs = Presentation(str(pptx_path))
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    slides = list(prs.slides)

    from pptx.oxml.ns import qn
    from copy import deepcopy

    # Transfer YAML title-slide notes to the first slide (pandoc does not do this)
    if title_notes and slides and not slides[0].has_notes_slide:
        notes_slide = slides[0].notes_slide  # creates notes slide if absent
        notes_slide.notes_text_frame.text = title_notes

    if not slide_specs:
        print("No slide specs found in markdown.")
        return

    # Build lookup: normalized title prefix → spec
    spec_lookup = {}
    for spec in slide_specs:
        key = normalize(spec["title"])
        spec_lookup[key] = spec

    # Identify parent slides (match to spec by title) and orphans
    orphan_indices = set()
    parent_map = {}  # slide_index → spec
    last_parent_idx = None  # track for orphan → parent note transfer
    orphan_to_parent = {}  # orphan_idx → parent_idx

    for i, slide in enumerate(slides):
        title = normalize(get_slide_title(slide))
        # Try matching to a spec
        matched = False
        if title:  # skip empty titles
            best_score = 0
            best_spec = None
            for key, spec in spec_lookup.items():
                # Score by longest common prefix
                min_len = min(len(title), len(key))
                common = 0
                for c1, c2 in zip(title, key):
                    if c1 == c2:
                        common += 1
                    else:
                        break
                threshold = max(5, min(20, int(min_len * 0.8)))
                if common >= threshold and common > best_score:
                    best_score = common
                    best_spec = spec
            if best_spec:
                parent_map[i] = best_spec
                matched = True
                last_parent_idx = i

        if not matched:
            # Never remove the first slide (title slide) or slides with notes
            if i == 0:
                last_parent_idx = 0
                continue
            # Check if this slide's title is an image alt text (pandoc split)
            # Pandoc may concatenate alt text + bullet text, so use startswith
            is_image_split = False
            if title:
                for alt in all_image_alts:
                    if title.startswith(alt):
                        is_image_split = True
                        break
            if is_image_split:
                # This is a pandoc-generated split — always an orphan
                orphan_indices.add(i)
                if last_parent_idx is not None:
                    orphan_to_parent[i] = last_parent_idx
                continue
            # Untitled slides are pandoc continuation splits — always merge
            has_title_ph = False
            for shape in slide.shapes:
                try:
                    if (shape.placeholder_format is not None
                            and shape.placeholder_format.idx == 0
                            and shape.text_frame.text.strip()):
                        has_title_ph = True
                        break
                except (ValueError, AttributeError):
                    pass
            if not has_title_ph and not is_image_split:
                orphan_indices.add(i)
                if last_parent_idx is not None:
                    orphan_to_parent[i] = last_parent_idx
                continue
            if has_real_content(slide):
                last_parent_idx = i
                continue
            # This is an orphan
            orphan_indices.add(i)
            if last_parent_idx is not None:
                orphan_to_parent[i] = last_parent_idx

    # Transfer notes AND text content from orphan slides to their parent
    for orphan_idx, parent_idx in orphan_to_parent.items():
        orphan = slides[orphan_idx]
        parent = slides[parent_idx]
        # Transfer notes
        if orphan.has_notes_slide:
            orphan_notes = orphan.notes_slide.notes_text_frame.text.strip()
            if orphan_notes and not parent.has_notes_slide:
                parent_notes_slide = parent.notes_slide  # creates if needed
                parent_notes_slide.notes_text_frame.text = orphan_notes
        # Transfer text content from orphan body to parent body (text frames only)
        orphan_body = get_body_placeholder(orphan)
        parent_body = get_body_placeholder(parent)
        if orphan_body and parent_body:
            if not orphan_body.has_text_frame or not parent_body.has_text_frame:
                # Table or graphic placeholders — skip text transfer
                pass
            else:
                orphan_text = orphan_body.text_frame.text.strip()
                parent_text = parent_body.text_frame.text.strip()
                if orphan_text and not parent_text:
                    # Copy paragraph-by-paragraph to preserve formatting
                    src_txBody = orphan_body._element.find(qn("p:txBody"))
                    dst_txBody = parent_body._element.find(qn("p:txBody"))
                    if src_txBody is not None and dst_txBody is not None:
                        # Remove existing empty paragraphs from parent body
                        for p in list(dst_txBody.findall(qn("a:p"))):
                            dst_txBody.remove(p)
                        # Copy all paragraphs from orphan
                        for p in src_txBody.findall(qn("a:p")):
                            dst_txBody.append(deepcopy(p))
        # Transfer standalone table shapes from orphan to parent (multi-table slides).
        # When a markdown slide has two tables, pandoc splits the second onto a
        # continuation slide with no title.  We copy the table graphicFrame element
        # directly into the parent's shape tree.  layout_stacked_tables() will
        # reposition all tables into equal vertical bands later.
        orphan_tables = get_table_frames(orphan)
        if orphan_tables:
            spTree = parent._element.find(qn("p:cSld")).find(qn("p:spTree"))
            # Build set of existing shape IDs to avoid duplicates
            existing_ids: set[int] = set()
            for sp in spTree:
                for cNvPr in sp.findall(f".//{qn('p:cNvPr')}"):
                    try:
                        existing_ids.add(int(cNvPr.get("id", "0")))
                    except ValueError:
                        pass
            next_id = max(existing_ids) + 1 if existing_ids else 100
            for tbl_shape in orphan_tables:
                tbl_elem = deepcopy(tbl_shape._element)
                # Re-assign shape ID to avoid conflicts within the slide
                for cNvPr in tbl_elem.findall(f".//{qn('p:cNvPr')}"):
                    cNvPr.set("id", str(next_id))
                    next_id += 1
                spTree.append(tbl_elem)

    # Remove existing images from parent slides (will re-add from source)
    for idx, spec in parent_map.items():
        if not spec["images"]:
            continue
        slide = slides[idx]
        spec_alts = {normalize(img["alt"]) for img in spec["images"]}

        # Remove existing image shapes
        for img_shape in get_slide_images(slide):
            sp = img_shape._element
            sp.getparent().remove(sp)

        # Remove non-placeholder text shapes containing image alt text
        for sh in list(slide.shapes):
            try:
                if sh.placeholder_format is not None:
                    continue  # keep placeholders
            except (ValueError, AttributeError):
                pass
            if sh.has_text_frame:
                txt = normalize(sh.text_frame.text)
                if txt in spec_alts:
                    sh._element.getparent().remove(sh._element)

        # Transfer body content from orphan slides to parent
        # Pandoc puts the title on the parent but the body (bullets) on the orphan
        for orphan_idx, parent_idx in orphan_to_parent.items():
            if parent_idx != idx:
                continue
            orphan = slides[orphan_idx]
            orphan_body = get_body_placeholder(orphan)
            if orphan_body is None or not orphan_body.has_text_frame:
                continue

            # Clean alt text lines from orphan body
            orphan_txBody = orphan_body._element.find(qn("p:txBody"))
            if orphan_txBody is not None:
                for p in list(orphan_txBody.findall(qn("a:p"))):
                    p_text = normalize("".join(
                        r.text or "" for r in p.findall(f".//{qn('a:t')}")))
                    if p_text in spec_alts:
                        orphan_txBody.remove(p)

            remaining = orphan_body.text_frame.text.strip()
            if not remaining:
                continue

            # Check if parent already has a body placeholder
            parent_body = get_body_placeholder(slide)
            if parent_body is None:
                # Clone the orphan's body placeholder onto the parent slide
                slide._element.find(qn("p:cSld")).find(
                    qn("p:spTree")).append(deepcopy(orphan_body._element))
            else:
                # Merge: append orphan paragraphs to parent body
                dst_txBody = parent_body._element.find(qn("p:txBody"))
                src_txBody = orphan_body._element.find(qn("p:txBody"))
                if dst_txBody is not None and src_txBody is not None:
                    # Remove empty paragraphs from parent
                    for p in list(dst_txBody.findall(qn("a:p"))):
                        p_text = "".join(
                            r.text or "" for r in p.findall(f".//{qn('a:t')}")).strip()
                        if not p_text:
                            dst_txBody.remove(p)
                    for p in src_txBody.findall(qn("a:p")):
                        dst_txBody.append(deepcopy(p))

        # Decide layout based on whether slide now has text content
        image_paths = [img["path"] for img in spec["images"]]
        if slide_has_text_content(slide):
            layout_images_with_text(slide, image_paths, slide_w, slide_h)
        else:
            layout_images_on_slide(slide, image_paths, slide_w, slide_h)

    # Remove orphan slides
    if orphan_indices:
        sldIdLst = prs.slides._sldIdLst
        sldIds = list(sldIdLst)
        for idx in sorted(orphan_indices, reverse=True):
            if idx < len(sldIds):
                sldIdLst.remove(sldIds[idx])

    # Re-layout slides that now have multiple stacked tables (multi-table merge)
    for slide in prs.slides:
        layout_stacked_tables(slide, slide_w, slide_h)

    # Style all tables: navy header, alternating rows
    style_tables(prs)

    prs.save(str(pptx_path))

    # python-pptx only removes slides from sldIdLst; the XML files and rels
    # remain in the zip and cause PowerPoint to report "content has problems".
    # Strip the orphaned parts so the package is self-consistent.
    purged = purge_orphaned_slides(pptx_path)

    total = len(prs.slides)
    img_slides = sum(1 for spec in parent_map.values() if spec["images"])
    print(f"Post-processed: {len(orphan_indices)} orphan(s) removed, "
          f"{img_slides} slide(s) got images. Result: {total} slides. "
          f"Purged {purged} stale part(s). Saved {pptx_path}")


def purge_orphaned_slides(pptx_path: Path) -> int:
    """Remove slide XML files from the zip that are not referenced by sldIdLst.

    python-pptx's sldIdLst.remove() removes slides from the visible presentation
    sequence but leaves the underlying XML files and relationships inside the zip.
    PowerPoint validates package consistency and raises a "content has problems"
    error when it finds parts that are registered in relationships but absent from
    the slide list, or parts present in the zip with no referencing relationship.

    This function re-opens the saved zip, reconciles the slide list against the
    relationship map, and rewrites the zip with the orphaned parts excised.

    Returns the number of orphaned slide parts removed.
    """
    import io
    import zipfile as _zipfile
    from lxml import etree as _etree

    SLIDE_REL = ("http://schemas.openxmlformats.org/officeDocument/2006/"
                 "relationships/slide")
    NOTES_REL = ("http://schemas.openxmlformats.org/officeDocument/2006/"
                 "relationships/notesSlide")
    NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"

    raw = pptx_path.read_bytes()

    with _zipfile.ZipFile(io.BytesIO(raw)) as zin:
        names = set(zin.namelist())

        # --- which rIds are actually in sldIdLst? ---
        prs_xml = _etree.fromstring(zin.read("ppt/presentation.xml"))
        active_rids: set[str] = set()
        for el in prs_xml.iter(f"{{{NS_P}}}sldId"):
            rid = el.get(f"{{{NS_R}}}id")
            if rid:
                active_rids.add(rid)

        # --- full slide relationship map ---
        prs_rels_xml = _etree.fromstring(
            zin.read("ppt/_rels/presentation.xml.rels"))
        slide_rels: dict[str, str] = {}  # rId -> target (relative to ppt/)
        for rel in prs_rels_xml:
            if rel.get("Type") == SLIDE_REL:
                slide_rels[rel.get("Id")] = rel.get("Target")

        # --- orphaned rIds ---
        orphaned: dict[str, str] = {
            rid: tgt for rid, tgt in slide_rels.items()
            if rid not in active_rids
        }
        if not orphaned:
            return 0

        # --- collect every file to drop ---
        drop: set[str] = set()
        for rid, target in orphaned.items():
            # target is like "slides/slide18.xml" (relative to ppt/)
            slide_file = f"ppt/{target}"
            slide_name = target.split("/")[-1]          # slide18.xml
            slide_rels_file = f"ppt/slides/_rels/{slide_name}.rels"
            drop.add(slide_file)
            drop.add(slide_rels_file)

            # find the associated notes slide
            if slide_rels_file in names:
                srels = _etree.fromstring(zin.read(slide_rels_file))
                for r in srels:
                    if r.get("Type") == NOTES_REL:
                        notes_name = r.get("Target", "").split("/")[-1]
                        drop.add(f"ppt/notesSlides/{notes_name}")
                        drop.add(f"ppt/notesSlides/_rels/{notes_name}.rels")

        # --- remove orphaned rels from presentation.xml.rels ---
        for rel in list(prs_rels_xml):
            if rel.get("Id") in orphaned:
                prs_rels_xml.remove(rel)
        new_prs_rels = _etree.tostring(
            prs_rels_xml, xml_declaration=True,
            encoding="UTF-8", standalone=True)

        # --- remove Override entries from [Content_Types].xml ---
        ct_xml = _etree.fromstring(zin.read("[Content_Types].xml"))
        ct_ns = "http://schemas.openxmlformats.org/package/2006/content-types"
        for ovr in list(ct_xml.findall(f"{{{ct_ns}}}Override")):
            part_name = ovr.get("PartName", "").lstrip("/")
            if part_name in drop:
                ct_xml.remove(ovr)
        new_ct = _etree.tostring(
            ct_xml, xml_declaration=True,
            encoding="UTF-8", standalone=True)

        # --- rebuild zip ---
        buf = io.BytesIO()
        with _zipfile.ZipFile(buf, "w", _zipfile.ZIP_DEFLATED) as zout:
            for name in zin.namelist():
                if name in drop:
                    continue
                if name == "ppt/_rels/presentation.xml.rels":
                    zout.writestr(name, new_prs_rels)
                elif name == "[Content_Types].xml":
                    zout.writestr(name, new_ct)
                else:
                    zout.writestr(name, zin.read(name))

    pptx_path.write_bytes(buf.getvalue())
    return len(orphaned)


def main() -> int:
    args = sys.argv[1:]
    resource_path = None

    # Parse --resource-path
    filtered = []
    i = 0
    while i < len(args):
        if args[i] == "--resource-path" and i + 1 < len(args):
            resource_path = Path(args[i + 1])
            i += 2
        else:
            filtered.append(args[i])
            i += 1

    if len(filtered) < 2:
        print("usage: python scripts/postprocess_slides.py slides.pptx slides.md "
              "[--resource-path DIR]", file=sys.stderr)
        return 2

    pptx_path = Path(filtered[0])
    md_path = Path(filtered[1])

    if not pptx_path.exists():
        print(f"pptx not found: {pptx_path}", file=sys.stderr)
        return 1
    if not md_path.exists():
        print(f"markdown not found: {md_path}", file=sys.stderr)
        return 1

    if resource_path is None:
        resource_path = md_path.parent

    postprocess(pptx_path, md_path, resource_path)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
