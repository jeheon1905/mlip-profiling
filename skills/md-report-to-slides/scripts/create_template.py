#!/usr/bin/env python3
"""Generate a styled reference pptx template for pandoc slide rendering.

Usage:
    python scripts/create_template.py [output.pptx]

Extracts pandoc's default reference.pptx, then applies visual design elements
(colored header bars, accent lines, font styling) to the slide layouts.
The result is used with pandoc's --reference-doc flag.

Requires: python-pptx, pandoc
"""
from __future__ import annotations

import subprocess
import sys
import tempfile
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from lxml import etree
except ImportError:
    print("python-pptx and lxml are required: pip install python-pptx lxml",
          file=sys.stderr)
    sys.exit(1)

# ---------------------------------------------------------------------------
# Theme configuration
# ---------------------------------------------------------------------------
THEME = {
    # Colors
    "navy": RGBColor(0x1B, 0x2A, 0x4A),
    "dark_text": RGBColor(0x2D, 0x2D, 0x2D),
    "accent": RGBColor(0x26, 0x6D, 0xD3),
    "light_accent": RGBColor(0x3D, 0x8B, 0xF8),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "light_bg": RGBColor(0xF5, 0xF7, 0xFA),
    "subtle_grey": RGBColor(0xAA, 0xAA, 0xAA),
    # Fonts
    "title_font": "Segoe UI Semibold",
    "body_font": "Segoe UI",
    # Sizes
    "title_slide_title": Pt(32),
    "title_slide_subtitle": Pt(16),
    "content_title": Pt(22),
    "content_body": Pt(16),
    # Dimensions
    "header_height": Inches(1.1),
    "accent_line_height": Pt(3),
}


def extract_pandoc_reference() -> Path:
    """Extract pandoc's built-in reference.pptx to a temp file."""
    tmp = Path(tempfile.mktemp(suffix=".pptx"))
    result = subprocess.run(
        ["pandoc", "--print-default-data-file", "reference.pptx"],
        capture_output=True,
    )
    if result.returncode != 0:
        print("Failed to extract pandoc reference.pptx.", file=sys.stderr)
        sys.exit(1)
    tmp.write_bytes(result.stdout)
    return tmp


def rgb_hex(color: RGBColor) -> str:
    return f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"


def add_rect_to_layout(layout, left, top, width, height, color: RGBColor):
    """Add a filled rectangle to a slide layout via XML (LayoutShapes has no add_shape)."""
    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    nsp = "http://schemas.openxmlformats.org/presentationml/2006/main"
    nsr = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

    spTree = layout._element.find(f"{{{nsp}}}cSld/{{{nsp}}}spTree")
    if spTree is None:
        spTree = layout._element.find(f".//{{{nsp}}}spTree")
    if spTree is None:
        return

    # Generate a unique id
    existing_ids = [int(sp.find(f"{{{nsp}}}nvSpPr/{{{nsp}}}cNvPr").get("id", "0"))
                    for sp in spTree.findall(f"{{{nsp}}}sp")
                    if sp.find(f"{{{nsp}}}nvSpPr/{{{nsp}}}cNvPr") is not None]
    new_id = max(existing_ids, default=100) + 1

    sp_xml = f'''<p:sp xmlns:p="{nsp}" xmlns:a="{ns}" xmlns:r="{nsr}">
      <p:nvSpPr>
        <p:cNvPr id="{new_id}" name="Rect{new_id}"/>
        <p:cNvSpPr/>
        <p:nvPr/>
      </p:nvSpPr>
      <p:spPr>
        <a:xfrm>
          <a:off x="{int(left)}" y="{int(top)}"/>
          <a:ext cx="{int(width)}" cy="{int(height)}"/>
        </a:xfrm>
        <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
        <a:solidFill>
          <a:srgbClr val="{rgb_hex(color)}"/>
        </a:solidFill>
        <a:ln w="0">
          <a:noFill/>
        </a:ln>
      </p:spPr>
    </p:sp>'''

    sp_elem = etree.fromstring(sp_xml)
    # Insert at beginning so it's behind other shapes
    spTree.insert(0, sp_elem)


def style_defRPr(pPr, font_name: str, font_size_pt: float, color: RGBColor,
                 bold: bool = False):
    """Set default run properties on a paragraph's pPr element."""
    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    defRPr = pPr.find(f"{{{ns}}}defRPr")
    if defRPr is None:
        defRPr = etree.SubElement(pPr, f"{{{ns}}}defRPr")

    defRPr.set("sz", str(int(font_size_pt * 100)))
    defRPr.set("b", "1" if bold else "0")

    # Remove old font/color children
    for child in list(defRPr):
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag in ("latin", "solidFill", "ea", "cs"):
            defRPr.remove(child)

    latin = etree.SubElement(defRPr, f"{{{ns}}}latin")
    latin.set("typeface", font_name)
    ea = etree.SubElement(defRPr, f"{{{ns}}}ea")
    ea.set("typeface", font_name)

    solidFill = etree.SubElement(defRPr, f"{{{ns}}}solidFill")
    srgbClr = etree.SubElement(solidFill, f"{{{ns}}}srgbClr")
    srgbClr.set("val", rgb_hex(color))


def style_paragraph(para, font_name: str, font_size, font_color: RGBColor,
                    bold: bool = False):
    """Apply font styling to a paragraph and its runs."""
    pPr = para._p.get_or_add_pPr()
    style_defRPr(pPr, font_name, font_size.pt, font_color, bold)

    for run in para.runs:
        run.font.name = font_name
        run.font.size = font_size
        run.font.color.rgb = font_color
        run.font.bold = bold


def style_title_slide_layout(layout, slide_w, slide_h):
    """Style the Title Slide layout with dark background and accent."""
    t = THEME

    # Full-slide dark background
    bg = layout.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = t["navy"]

    # Accent line near bottom third
    accent_y = int(slide_h * 0.62)
    add_rect_to_layout(layout, Inches(0.8), accent_y,
             Inches(2.5), t["accent_line_height"], t["accent"])

    # Style placeholders
    for ph in layout.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:  # Title
            ph.left = Inches(0.8)
            ph.top = int(slide_h * 0.28)
            ph.width = int(slide_w - Inches(1.6))
            ph.height = Inches(1.5)
            for para in ph.text_frame.paragraphs:
                style_paragraph(para, t["title_font"], t["title_slide_title"],
                                t["white"], bold=True)
                para.alignment = PP_ALIGN.LEFT
        elif idx == 1:  # Subtitle
            ph.left = Inches(0.8)
            ph.top = int(slide_h * 0.65)
            ph.width = int(slide_w - Inches(1.6))
            ph.height = Inches(1.0)
            for para in ph.text_frame.paragraphs:
                style_paragraph(para, t["body_font"], t["title_slide_subtitle"],
                                t["subtle_grey"])
                para.alignment = PP_ALIGN.LEFT


def style_content_layout(layout, slide_w, slide_h):
    """Style content slide layout with colored header bar and accent line."""
    t = THEME

    # White background
    bg = layout.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = t["white"]

    # Header bar (navy)
    add_rect_to_layout(layout, 0, 0, slide_w, t["header_height"], t["navy"])

    # Accent line under header
    add_rect_to_layout(layout, 0, t["header_height"],
             slide_w, t["accent_line_height"], t["accent"])

    # Style placeholders
    for ph in layout.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:  # Title — sits inside the header bar
            ph.left = Inches(0.6)
            ph.top = Inches(0.15)
            ph.width = int(slide_w - Inches(1.2))
            ph.height = Inches(0.8)
            ph.text_frame.word_wrap = True
            for para in ph.text_frame.paragraphs:
                style_paragraph(para, t["title_font"], t["content_title"],
                                t["white"], bold=True)
                para.alignment = PP_ALIGN.LEFT
        elif idx == 1:  # Body — below header
            ph.left = Inches(0.6)
            ph.top = Inches(1.3)
            ph.width = int(slide_w - Inches(1.2))
            ph.height = int(slide_h - Inches(1.6))
            ph.text_frame.word_wrap = True
            for para in ph.text_frame.paragraphs:
                style_paragraph(para, t["body_font"], t["content_body"],
                                t["dark_text"])
                para.alignment = PP_ALIGN.LEFT


def style_two_content_layout(layout, slide_w, slide_h):
    """Style two-content (comparison) layout — same header as content."""
    t = THEME

    bg = layout.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = t["white"]

    add_rect_to_layout(layout, 0, 0, slide_w, t["header_height"], t["navy"])
    add_rect_to_layout(layout, 0, t["header_height"],
             slide_w, t["accent_line_height"], t["accent"])

    for ph in layout.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            ph.left = Inches(0.6)
            ph.top = Inches(0.15)
            ph.width = int(slide_w - Inches(1.2))
            ph.height = Inches(0.8)
            for para in ph.text_frame.paragraphs:
                style_paragraph(para, t["title_font"], t["content_title"],
                                t["white"], bold=True)


def create_template(output_path: Path):
    ref_path = extract_pandoc_reference()
    prs = Presentation(str(ref_path))
    ref_path.unlink(missing_ok=True)

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # Style each slide layout by name
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            name = layout.name.lower()
            if "title slide" in name or name == "title slide":
                style_title_slide_layout(layout, slide_w, slide_h)
            elif "two content" in name:
                style_two_content_layout(layout, slide_w, slide_h)
            elif "content" in name or "blank" not in name:
                style_content_layout(layout, slide_w, slide_h)

    # Remove existing slides
    xml_slides = prs.slides._sldIdLst
    for sld in list(xml_slides):
        xml_slides.remove(sld)

    prs.save(str(output_path))
    print(f"Template created: {output_path}")
    print(f"Usage: pandoc slides.md -o slides.pptx --reference-doc={output_path}")


def main() -> int:
    output = Path(sys.argv[1]) if len(sys.argv) > 1 else Path("template.pptx")
    create_template(output)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
